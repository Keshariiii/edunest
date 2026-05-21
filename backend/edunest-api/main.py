import os
import json
import base64
import asyncio
import traceback
import httpx
from io import BytesIO
from typing import List, Any, Optional
from fastapi import FastAPI, File, Form, UploadFile, HTTPException
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

# Optional: docx/pptx support
try:
    # pyrefly: ignore [missing-import]
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("WARNING: python-docx not installed. .docx uploads will be rejected.")

try:
    # pyrefly: ignore [missing-import]
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("WARNING: python-pptx not installed. .pptx uploads will be rejected.")

# Initialize FastAPI app for EduNest
app = FastAPI()

# Enable CORS for local & network development (phone, LAN devices, etc.)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── Database initialization ─────────────────────────────────────────────────────
from database import engine, Base, SessionLocal
from models import User, PasswordResetToken, RefreshToken, FailedLoginAttempt
Base.metadata.create_all(bind=engine)


def lowercase_existing_emails(session_factory=None):
    """One-time startup migration: normalise every stored email to lowercase.

    If two accounts share the same email after lowercasing (e.g. User@x.com and
    user@x.com), the *newest* account (highest id) keeps the clean lowercase
    email, while older duplicates are renamed to ``{username}_dup_{id}@domain``
    so the UNIQUE constraint is not violated.

    Args:
        session_factory: Optional SQLAlchemy sessionmaker. Defaults to the
            module-level ``SessionLocal`` when None (production path).
    """
    if session_factory is None:
        session_factory = SessionLocal
    db = session_factory()
    try:
        users = db.query(User).order_by(User.id.desc()).all()
        seen: dict[str, int] = {}          # lowercase email → user.id that owns it
        changed = 0
        for user in users:
            lower = user.email.lower()
            if lower in seen:
                # This is an older duplicate — rename it
                local, _, domain = lower.partition("@")
                user.email = f"{user.username}_dup_{user.id}@{domain}"
                changed += 1
                print(f"[EMAIL MIGRATION] Renamed duplicate id={user.id} to {user.email}")
            else:
                seen[lower] = user.id
                if user.email != lower:
                    user.email = lower
                    changed += 1
        # Also lowercase FailedLoginAttempt emails
        attempts = db.query(FailedLoginAttempt).all()
        for attempt in attempts:
            low = attempt.email.lower()
            if attempt.email != low:
                attempt.email = low
                changed += 1
        if changed:
            db.commit()
            print(f"[EMAIL MIGRATION] Lowercased {changed} record(s).")
        else:
            print("[EMAIL MIGRATION] All emails already lowercase — nothing to do.")
    except Exception as e:
        db.rollback()
        print(f"[EMAIL MIGRATION] Error: {e}")
    finally:
        db.close()


lowercase_existing_emails()

# ── Auth routes ─────────────────────────────────────────────────────────────────
from auth_routes import router as auth_router, user_router
app.include_router(auth_router)
app.include_router(user_router)

# Fix Swagger UI file upload: patch openapi schema to use format: binary
from fastapi.openapi.utils import get_openapi

def custom_openapi():
    if app.openapi_schema:
        return app.openapi_schema
    schema = get_openapi(title=app.title, version=app.version, routes=app.routes)
    for path_data in schema.get("paths", {}).values():
        for operation in path_data.values():
            body = operation.get("requestBody", {})
            content = body.get("content", {})
            form = content.get("multipart/form-data", {})
            schema_ref = form.get("schema", {})
            ref = schema_ref.get("$ref", "")
            if ref:
                def_name = ref.split("/")[-1]
                body_schema = schema.get("components", {}).get("schemas", {}).get(def_name, {})
                props = body_schema.get("properties", {})
                if "files" in props and props["files"].get("type") == "array":
                    props["files"]["items"] = {"type": "string", "format": "binary"}
    app.openapi_schema = schema
    return app.openapi_schema

app.openapi = custom_openapi

from fastapi.responses import JSONResponse
from fastapi import Request
import traceback
@app.exception_handler(Exception)
async def generic_exception_handler(request: Request, exc: Exception):
    return JSONResponse(
        status_code=500,
        content={"detail": "Internal Server Error", "traceback": traceback.format_exc()}
    )

# ─────────────────────────────────────────────────────────────────────────────
# GET /api/health — lightweight ping for frontend status indicator
# ─────────────────────────────────────────────────────────────────────────────
@app.get("/api/health")
async def health_check():
    return {"status": "ok"}

# API key is loaded dynamically on every request so swapping the key in .env
# takes effect immediately — no server restart required.
try:
    from dotenv import load_dotenv
except ImportError:
    load_dotenv = None  # type: ignore

def get_api_key() -> str:
    """Re-reads GEMINI_API_KEY from .env on every call so key changes are instant."""
    if load_dotenv is not None:
        from pathlib import Path
        env_path = Path(__file__).parent / ".env"
        load_dotenv(dotenv_path=env_path, override=True)
    key = os.getenv("GEMINI_API_KEY")
    if not key:
        raise ValueError("No API key found. Please set GEMINI_API_KEY in your .env file.")
    return key

GEMINI_MODEL = "gemini-2.5-flash-lite"
GEMINI_API_URL = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent"

# ─────────────────────────────────────────────────────────────────────────────
# Core helper: calls Gemini REST API directly with a 5-minute timeout.
# This completely bypasses the google-genai SDK which has a known bug where
# httpx timeouts are silently ignored in async mode.
# ─────────────────────────────────────────────────────────────────────────────
async def call_gemini(parts: list, generation_config: dict) -> str:
    """
    Calls the Gemini REST API directly via httpx with a 5-minute timeout.
    Re-reads the API key from .env on every call — swap the key anytime without restart.
    Auto-retries on 503 (server overloaded) and 429 (rate limit) with exponential backoff.
    Raises HTTP 429 HTTPException on quota exhaustion so the frontend receives a clean error.
    """
    api_key = get_api_key()  # ← fresh key read on every request
    payload = {
        "contents": [{"parts": parts}],
        "generationConfig": generation_config,
    }

    # Retry delays for 429/503: 20s on first retry, 45s on second
    # These are longer than before because 2.5-flash-lite has tighter per-minute limits.
    retry_delays = [20, 45]
    max_attempts = 3

    for attempt in range(max_attempts):
        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(
                f"{GEMINI_API_URL}?key={api_key}",
                json=payload,
            )

        # 503 — Google server overloaded
        if response.status_code == 503:
            wait = retry_delays[min(attempt, len(retry_delays) - 1)]
            print(f"503 overloaded — retrying in {wait}s (attempt {attempt+1}/{max_attempts})...")
            if attempt < max_attempts - 1:
                await asyncio.sleep(wait)
                continue
            raise HTTPException(
                status_code=503,
                detail="Gemini API is overloaded right now. Please wait a few seconds and try again."
            )

        # 429 — rate limit / quota exceeded
        if response.status_code == 429:
            # Honour Retry-After header if Google sends one, else use our default
            retry_after_hdr = response.headers.get("Retry-After", "")
            try:
                wait = int(retry_after_hdr)
            except (ValueError, TypeError):
                wait = retry_delays[min(attempt, len(retry_delays) - 1)]
            print(f"429 rate-limited — retrying in {wait}s (attempt {attempt+1}/{max_attempts})...")
            if attempt < max_attempts - 1:
                await asyncio.sleep(wait)
                continue
            # All retries exhausted — tell the frontend exactly what happened
            raise HTTPException(
                status_code=429,
                detail=(
                    "Gemini API rate limit hit. "
                    "gemini-2.5-flash-lite has a low free-tier quota. "
                    "Please wait 60 seconds and try again, or reduce your file size."
                )
            )

        # All other errors (400, 500, etc.) — raise immediately
        response.raise_for_status()
        break

    result = response.json()
    try:
        text = result["candidates"][0]["content"]["parts"][0]["text"]
    except (KeyError, IndexError) as e:
        raise RuntimeError(f"Unexpected Gemini response format: {result}") from e
    return text


def parse_json_response(raw: str, raise_on_error: bool = True) -> dict:
    """Strips optional markdown fences and parses JSON.
    
    Args:
        raw: The raw string response from Gemini.
        raise_on_error: If True (default), raises RuntimeError on JSON parse failure.
                        If False, returns a graceful fallback dict (used by study-material endpoint).
    """
    text = raw.strip()
    if text.startswith("```json"):
        text = text[7:]
    elif text.startswith("```"):
        text = text[3:]
    if text.endswith("```"):
        text = text[:-3]
    try:
        return json.loads(text.strip())
    except json.JSONDecodeError as e:
        print(f"JSON Parsing Error: {e}")
        if raise_on_error:
            raise RuntimeError(
                f"The AI returned invalid JSON that could not be parsed. "
                f"Error: {e}. First 300 chars of response: {text[:300]}"
            ) from e
        # Graceful fallback for study-material endpoint only
        return {
            "formulas": [],
            "short_notes": ["Analysis completed, but the AI response contained invalid formatting.", text.strip()],
            "flashcards": []
        }


import io
from PIL import Image

def extract_text_from_docx(data: bytes) -> str:
    """Extract plain text from a .docx file."""
    try:
        doc = DocxDocument(BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        raise RuntimeError(f"Failed to parse .docx: {e}")

def extract_text_from_pptx(data: bytes) -> str:
    """Extract plain text from a .pptx file."""
    try:
        prs = PptxPresentation(BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        raise RuntimeError(f"Failed to parse .pptx: {e}")


def encode_file(data: bytes, mime_type: str) -> dict:
    """
    Encodes raw file bytes as a Gemini inlineData part.
    Images are aggressively compressed to 512x512 / JPEG-55 to minimise token
    usage — the #1 cause of 429 quota exhaustion on 2.5-flash-lite's free tier.
    """
    if mime_type.startswith("image/"):
        try:
            img = Image.open(io.BytesIO(data))
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            # Keep max dimension at 512 px — still legible for text/equations,
            # but uses ~4× fewer tokens than 1024 px.
            img.thumbnail((512, 512))

            output = io.BytesIO()
            img.save(output, format="JPEG", quality=55)
            compressed = output.getvalue()

            # Safety: never use a compressed image larger than the original
            if len(compressed) < len(data):
                data = compressed
            mime_type = "image/jpeg"
            print(f"Image compressed to {len(data)//1024} KB")
        except Exception as e:
            print(f"Warning: Failed to compress image: {e}")

    return {
        "inlineData": {
            "mimeType": mime_type,
            "data": base64.b64encode(data).decode("utf-8"),
        }
    }


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/generate-study-material
# ─────────────────────────────────────────────────────────────────────────────
@app.post("/api/generate-study-material")
async def generate_study_material(
    files: List[UploadFile] = File(...),
    subject: str = Form(...),
):
    try:
        # Build file parts (inline base64), with text extraction for office formats
        file_parts = []
        for file in files:
            file_data = await file.read()
            fname = (file.filename or "").lower()
            ct = file.content_type or ""
            if fname.endswith(".docx") or ct == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                if not DOCX_AVAILABLE:
                    raise HTTPException(status_code=400, detail="python-docx not installed on server.")
                extracted = extract_text_from_docx(file_data)
                file_parts.append({"text": f"[DOCUMENT CONTENT]:\n{extracted}"})
            elif fname.endswith(".pptx") or ct == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                if not PPTX_AVAILABLE:
                    raise HTTPException(status_code=400, detail="python-pptx not installed on server.")
                extracted = extract_text_from_pptx(file_data)
                file_parts.append({"text": f"[SLIDE CONTENT]:\n{extracted}"})
            else:
                file_parts.append(encode_file(file_data, ct or "application/octet-stream"))

        prompt_text = f"""
        ACT AS AN EXPERT TUTOR FOR {subject.upper()}.
        ANALYZE THE PROVIDED DOCUMENT TEXT AND/OR IMAGES AND EXTRACT HIGH-YIELD STUDY MATERIAL.

        CRITICAL INSTRUCTION: Aggressively identify any mathematical relationships, equations, or theorems to populate the Formulas tab, even if they aren't written in standard LaTeX format.

        !!! CRITICAL JSON RULE !!!
        You MUST double-escape all backslashes in your LaTeX (e.g., use \\\\frac instead of \\frac, \\\\sum instead of \\sum). This is strictly required for valid JSON.

        TASK & STRICT LIMITS:
        0. Metadata: Identify the chapter/topic name from this document. List the 3-6 key topics covered.
        1. Formulas: Extract the 5 to 10 most important mathematical/scientific formulas from this document. Convert to LaTeX.
        2. Short Notes: Generate 5 to 10 highly condensed, core conceptual points from this document.
           Return clean, plain text strings for short_notes. DO NOT start strings with dashes (-), bullet points, or asterisks (*).
        3. Flashcards: Create exactly 5 to 10 active-recall flashcards covering high-yield definitions and concepts.
           !!! MATH IN FLASHCARDS RULE !!!
           - ALL mathematical expressions, equations, variables, and symbols MUST be wrapped in LaTeX inline delimiters: $...$
           - Example: write "$x^3 - 6x + 11x - 6$" NOT "x^3 - 6x + 11x - 6"
           - Example: write "$a^m \\\\times a^n = a^{{m+n}}$" NOT "a^m \\times a^n = a^{{m+n}}"
           - Double-escape ALL backslashes inside $...$: use \\\\frac, \\\\sqrt, \\\\times, \\\\sum, etc.
           - Plain text parts of the flashcard (non-math) should remain as normal text.

        FORMAT EXACTLY AS THIS JSON STRUCTURE:
        {{
          "metadata": {{ "chapter_title": "string", "topics_covered": ["topic1", "topic2"] }},
          "formulas": [{{ "name": "string", "equation": "latex_string" }}],
          "short_notes": ["clean string without bullet characters"],
          "flashcards": [{{ "front": "string", "back": "string" }}]
        }}
        """

        parts = [{"text": prompt_text}] + file_parts

        generation_config = {
            "responseMimeType": "application/json",
            "temperature": 0.3,
        }

        max_retries = 3
        raw_response = None
        for attempt in range(max_retries):
            try:
                print(f"Study material attempt {attempt + 1}/{max_retries}...")
                raw_response = await call_gemini(parts, generation_config)
                break
            except HTTPException:
                # Propagate clean 429/503 from call_gemini directly to the client
                raise
            except httpx.TimeoutException as e:
                print(f"Timeout on attempt {attempt + 1}: {e}")
                if attempt < max_retries - 1:
                    await asyncio.sleep(2)
                else:
                    raise RuntimeError("Gemini API timed out after 5 minutes. Please try a smaller file.") from e

        data = parse_json_response(raw_response, raise_on_error=False)
        return {
            "metadata":    data.get("metadata", {"chapter_title": "", "topics_covered": []}),
            "formulas":    data.get("formulas", []),
            "short_notes": data.get("short_notes", []),
            "flashcards":  data.get("flashcards", []),
        }

    except HTTPException:
        raise  # pass 429/503 straight through
    except Exception as e:
        traceback.print_exc()
        print(f"STUDY MATERIAL ERROR: {e}")
        return JSONResponse(
            status_code=500,
            content={
                "error": True,
                "message": "Failed to process the document. Try a clearer photo or PDF.",
                "details": str(e),
            },
        )


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/generate-quiz
# ─────────────────────────────────────────────────────────────────────────────
@app.post("/api/generate-quiz")
async def generate_quiz(
    files: List[UploadFile] = File(...),
    subject: str = Form(...),
    exam_type: str = Form(default="JEE Main"),
    difficulty: str = Form(default="Medium"),
    num_questions: int = Form(default=10),
):
    try:
        # Encode files once, extracting text for office formats if needed
        file_parts = []
        for file in files:
            file_data = await file.read()
            fname = (file.filename or "").lower()
            ct = file.content_type or ""
            if fname.endswith(".docx") or ct == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                if not DOCX_AVAILABLE:
                    raise HTTPException(status_code=400, detail="python-docx not installed on server.")
                extracted = extract_text_from_docx(file_data)
                file_parts.append({"text": f"[DOCUMENT CONTENT]:\n{extracted}"})
            elif fname.endswith(".pptx") or ct == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                if not PPTX_AVAILABLE:
                    raise HTTPException(status_code=400, detail="python-pptx not installed on server.")
                extracted = extract_text_from_pptx(file_data)
                file_parts.append({"text": f"[SLIDE CONTENT]:\n{extracted}"})
            else:
                file_parts.append(encode_file(file_data, ct or "application/octet-stream"))

        # Use batch_size=10 so a typical 10-question quiz costs 1 API call, not 4.
        # Previously batch_size=3 caused 10 questions → 4 Gemini calls (3+3+3+1).
        batch_size = 10
        all_mcqs: List[dict[str, Any]] = []

        batches = [batch_size] * (num_questions // batch_size)
        if num_questions % batch_size != 0:
            batches.append(num_questions % batch_size)

        generation_config = {
            "responseMimeType": "application/json",
            "temperature": 0.4,
        }

        for i, num_in_batch in enumerate(batches):
            print(f"--- BATCH {i+1}/{len(batches)} ({num_in_batch} questions) ---")

            prompt_text = f"""
            Based on the provided notes, create exactly {num_in_batch} {exam_type} quiz questions for {subject}.
            Difficulty level: {difficulty}.

            !!! CRITICAL MATH FORMATTING RULE !!!
            Ensure all mathematical symbols, Greek letters, and units (like Omega) are enclosed in single dollar signs for inline math (e.g., $5 \\Omega$) or double dollar signs for block equations.
            BAD: Calculate mole fraction (x) using Henry's Law (p = K_H * x)
            GOOD: Calculate mole fraction ($x$) using Henry's Law ($p = K_H \\cdot x$)

            Whenever you output a Greek letter (such as \\alpha, \\beta, \\gamma, \\Delta, \\Omega, \\varepsilon, etc.), you MUST wrap it in a LaTeX textcolor command using the exact hex code #ff453a.
            BAD: Calculate the resistance using $R = \\rho \\frac{{L}}{{A}}$ and $\\Omega$
            GOOD: Calculate the resistance using $R = \\textcolor{{#ff453a}}{{\\rho}} \\frac{{L}}{{A}}$ and $\\textcolor{{#ff453a}}{{\\Omega}}$

            DO NOT repeat these specific concepts: {[q.get('question', '')[:50] for q in all_mcqs]}
            Keep explanations concise (maximum 4 steps).

            Return ONLY this JSON structure:
            {{ "mcqs": [{{
                "question": "string containing $math$",
                "options": ["string containing $math$"],
                "correct_indices": [int],
                "explanation": "string containing $math$",
                "type": "mcq"
            }}] }}
            """

            parts = [{"text": prompt_text}] + file_parts
            batch_success = False

            for attempt in range(3):
                raw_response = None
                try:
                    print(f"  Attempt {attempt + 1}/3...")
                    raw_response = await call_gemini(parts, generation_config)
                    batch_data = parse_json_response(raw_response)
                    if "mcqs" in batch_data:
                        all_mcqs.extend(batch_data["mcqs"])
                        print(f"  Added {len(batch_data['mcqs'])} questions.")
                        batch_success = True
                        break
                except HTTPException as e:
                    # 429 / 503 raised by call_gemini — retry with backoff
                    print(f"  Batch {i+1} HTTP {e.status_code} on attempt {attempt+1}: {e.detail}")
                    if e.status_code in (429, 503) and attempt < 2:
                        wait = 30 if e.status_code == 429 else 10
                        print(f"  Waiting {wait}s before retry...")
                        await asyncio.sleep(wait)
                    else:
                        print(f"  Batch {i+1} failed permanently.")
                except httpx.HTTPStatusError as e:
                    status = e.response.status_code
                    print(f"  Batch {i+1} API error {status} on attempt {attempt+1}")
                    if status == 429 and attempt < 2:
                        await asyncio.sleep(30)
                    else:
                        print(f"  Batch {i+1} failed.")
                except Exception as batch_err:
                    print(f"  Batch {i+1} error on attempt {attempt+1}: {batch_err}")
                    if raw_response:
                        print(f"  Raw: {raw_response[:300]}")
                    if attempt < 2:
                        await asyncio.sleep(15)

            if not batch_success:
                print(f"Stopping — batch {i+1} failed all retries.")
                break

            if (i < len(batches) - 1):
                print(f"Sleeping 5s between batches...")
                await asyncio.sleep(5)

        if len(all_mcqs) == 0:
            raise HTTPException(
                status_code=422,
                detail="The AI could not extract readable text from this document. Please try a clearer PDF.",
            )

        return {"mcqs": all_mcqs[:num_questions]}

    except HTTPException as he:
        raise he
    except Exception as e:
        traceback.print_exc()
        print(f"QUIZ SETUP ERROR: {e}")
        raise HTTPException(status_code=500, detail=f"Backend Error: {str(e)}")


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/evaluate-quiz
# ─────────────────────────────────────────────────────────────────────────────
class QuestionEvaluationItem(BaseModel):
    correct_indices: List[int]
    selected_indices: Optional[List[int]] = None

class QuizEvaluationRequest(BaseModel):
    questions: List[QuestionEvaluationItem]

@app.post("/api/evaluate-quiz")
async def evaluate_quiz(request: QuizEvaluationRequest):
    total_score = 0
    max_possible_score = 0
    individual_evaluations = []

    for q in request.questions:
        max_possible_score += 1
        if not q.selected_indices:
            marks_awarded = 0
            is_correct = None
        else:
            is_correct = sorted(q.selected_indices) == sorted(q.correct_indices)
            marks_awarded = 1 if is_correct else 0

        total_score += marks_awarded
        individual_evaluations.append({
            "is_correct": is_correct,
            "marks_awarded": marks_awarded,
            "correct_indices": q.correct_indices,
            "selected_indices": q.selected_indices if q.selected_indices else [],
        })

    return {
        "total_score": total_score,
        "max_possible_score": max_possible_score,
        "evaluations": individual_evaluations,
    }


# ─────────────────────────────────────────────────────────────────────────────
# AI DOUBT SOLVER — Merged Chatbot Backend
# Upload PDFs, index them with FAISS + Gemini embeddings, chat with context.
# ─────────────────────────────────────────────────────────────────────────────
import hashlib
import re
from pathlib import Path as PathLib

# Directories for doubt solver data
DOUBT_DATA_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "doubt_solver_data")
DOUBT_UPLOAD_DIR = os.path.join(DOUBT_DATA_DIR, "uploads")
DOUBT_INDEX_DIR = os.path.join(DOUBT_DATA_DIR, "index")
os.makedirs(DOUBT_UPLOAD_DIR, exist_ok=True)
os.makedirs(DOUBT_INDEX_DIR, exist_ok=True)

# In-memory session chat histories
_chat_sessions: dict[str, list[dict]] = {}

# Simple text chunking (no LangChain dependency needed)
def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 150) -> list[str]:
    """Split text into overlapping chunks."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap
    return [c.strip() for c in chunks if c.strip()]

# Extract text from PDF using PyPDF2 (already available via pypdf)
def extract_pdf_text(file_data: bytes) -> str:
    """Extract text from a PDF file."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(BytesIO(file_data))
        pages_text = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages_text.append(text)
        return "\n\n".join(pages_text)
    except ImportError:
        raise RuntimeError("pypdf not installed. Run: pip install pypdf")
    except Exception as e:
        raise RuntimeError(f"Failed to extract PDF text: {e}")



# ── Session-Isolated Chunk Store ─────────────────────────────────────────────
# Each chat session gets its own vector index and upload directory so that
# uploading files in one session does not bleed into another.

def _safe_session_id(session_id: str) -> str:
    """Sanitize session_id for use as a filesystem path component."""
    safe = "".join(c for c in session_id if c.isalnum() or c in ("-", "_")).strip()
    return safe or "global"

def _session_vector_file(session_id: str) -> str:
    return os.path.join(DOUBT_INDEX_DIR, f"vectors_{_safe_session_id(session_id)}.json")

def _session_upload_dir(session_id: str) -> str:
    d = os.path.join(DOUBT_UPLOAD_DIR, _safe_session_id(session_id))
    os.makedirs(d, exist_ok=True)
    return d

def load_vector_store(session_id: str = "global") -> dict:
    """Load the chunk store from disk for a given session."""
    vector_file = _session_vector_file(session_id)
    if os.path.exists(vector_file):
        with open(vector_file, "r", encoding="utf-8") as f:
            data = json.load(f)
            if "chunks" not in data:
                data["chunks"] = []
            if "metadata" not in data:
                data["metadata"] = []
            return data
    return {"chunks": [], "metadata": []}

def save_vector_store(store: dict, session_id: str = "global"):
    """Save the chunk store to disk for a given session."""
    vector_file = _session_vector_file(session_id)
    clean = {"chunks": store.get("chunks", []), "metadata": store.get("metadata", [])}
    with open(vector_file, "w", encoding="utf-8") as f:
        json.dump(clean, f)

def get_all_context(store: dict, max_chars: int = 80000) -> tuple[str, list[str]]:
    """Build a single context string from all stored chunks (up to max_chars).
    Returns (context_text, list_of_source_names)."""
    if not store["chunks"]:
        return "", []
    context_parts = []
    sources = set()
    total = 0
    for i, chunk in enumerate(store["chunks"]):
        if total + len(chunk) > max_chars:
            break
        context_parts.append(chunk)
        if i < len(store.get("metadata", [])):
            src = store["metadata"][i].get("source", "")
            if src:
                sources.add(src)
        total += len(chunk)
    return "\n\n---\n\n".join(context_parts), list(sources)



# Pydantic models for doubt solver
class DoubtChatRequest(BaseModel):
    session_id: str
    question: str

class DoubtChatResponse(BaseModel):
    answer: str
    sources: list[str] = []


# Image extensions supported for OCR via Gemini Vision
_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tiff", ".tif"}

async def extract_text_from_image(file_data: bytes, mime_type: str) -> str:
    """Use Gemini Vision to extract / transcribe ALL text from an image,
    including handwritten notes, diagrams with labels, and printed text."""
    image_part = encode_file(file_data, mime_type)
    prompt = {
        "text": (
            "You are an expert OCR and handwriting recognition system. "
            "Carefully examine this image and transcribe ALL text you can see — "
            "printed text, handwritten notes, equations, labels on diagrams, "
            "annotations, headings, bullet points, and any other written content. "
            "Preserve the original structure as much as possible (headings, lists, paragraphs). "
            "For mathematical equations, write them in a readable text format. "
            "If the image contains a diagram, briefly describe it and include any labels. "
            "Output ONLY the transcribed text, nothing else."
        )
    }
    generation_config = {"temperature": 0.1, "maxOutputTokens": 4096}
    return await call_gemini([prompt, image_part], generation_config)


@app.post("/api/doubt-solver/upload")
async def doubt_solver_upload(
    files: List[UploadFile] = File(...),
    session_id: str = "global",
):
    """
    Upload PDFs, documents, or images to the AI Chatbot.
    Extracts text (using OCR for images / handwritten notes), chunks it,
    and stores in a session-isolated chunk index.
    """
    if not files:
        raise HTTPException(status_code=400, detail="No files provided.")

    all_chunks = []
    all_metadata = []
    saved_files = []
    upload_dir = _session_upload_dir(session_id)

    try:
        for file in files:
            fname = (file.filename or "unknown").lower()
            file_data = await file.read()
            ct = file.content_type or ""
            ext = os.path.splitext(fname)[1]

            # Extract text based on file type
            if fname.endswith(".pdf"):
                text = extract_pdf_text(file_data)
            elif fname.endswith(".docx"):
                if not DOCX_AVAILABLE:
                    raise HTTPException(status_code=400, detail="python-docx not installed on server.")
                text = extract_text_from_docx(file_data)
            elif fname.endswith(".pptx"):
                if not PPTX_AVAILABLE:
                    raise HTTPException(status_code=400, detail="python-pptx not installed on server.")
                text = extract_text_from_pptx(file_data)
            elif fname.endswith(".txt"):
                text = file_data.decode("utf-8", errors="ignore")
            elif ext in _IMAGE_EXTENSIONS or ct.startswith("image/"):
                # Use Gemini Vision to OCR the image (handles handwritten notes)
                mime = ct if ct.startswith("image/") else f"image/{ext.lstrip('.')}"
                if mime == "image/jpg":
                    mime = "image/jpeg"
                print(f"  OCR-ing image: {file.filename} ({mime})")
                text = await extract_text_from_image(file_data, mime)
            else:
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type: {fname}. Use PDF, DOCX, PPTX, TXT, or images (JPG/PNG/WEBP)."
                )

            if not text.strip():
                continue

            # Save file to session-specific upload directory
            save_path = os.path.join(upload_dir, file.filename or "unknown")
            with open(save_path, "wb") as f:
                f.write(file_data)
            saved_files.append(file.filename)

            # Chunk the text
            chunks = chunk_text(text)
            for chunk in chunks:
                all_chunks.append(chunk)
                all_metadata.append({"source": file.filename, "type": ext.lstrip(".") or "unknown"})

        if not all_chunks:
            raise HTTPException(status_code=400, detail="No readable text found in the uploaded files.")

        # Store chunks (no embedding model needed)
        print(f"Storing {len(all_chunks)} chunks for session {session_id}...")
        store = load_vector_store(session_id)
        store["chunks"].extend(all_chunks)
        store["metadata"].extend(all_metadata)
        save_vector_store(store, session_id)

        return {
            "message": f"Successfully processed {len(saved_files)} file(s).",
            "files": saved_files,
            "chunks_indexed": len(all_chunks),
            "total_chunks_in_store": len(store["chunks"])
        }

    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Upload error: {str(e)}")


@app.post("/api/doubt-solver/chat", response_model=DoubtChatResponse)
async def doubt_solver_chat(request: DoubtChatRequest):
    """
    Chat with the AI Chatbot. Loads session-specific uploaded documents as context,
    then uses Gemini to generate an answer with conversational memory.
    """
    if not request.question.strip():
        raise HTTPException(status_code=400, detail="Question cannot be empty.")

    try:
        session_id = request.session_id
        store = load_vector_store(session_id)
        sources = []

        # Get all stored context for this session
        context_text, sources = get_all_context(store)
        if context_text:
            context_text = f"\n\nSTUDY MATERIAL CONTEXT:\n{context_text}\n"

        # Build chat history
        if session_id not in _chat_sessions:
            _chat_sessions[session_id] = []

        history = _chat_sessions[session_id]

        # Build conversation history string
        history_text = ""
        if history:
            recent = history[-10:]  # last 10 messages
            for msg in recent:
                role = "Student" if msg["role"] == "user" else "Tutor"
                history_text += f"{role}: {msg['content']}\n"

        system_prompt = f"""You are EduNest AI Tutor — a brilliant, patient, and encouraging study companion. 
Your role is to help students understand concepts from their uploaded study materials.

INSTRUCTIONS:
- If context from study materials is provided, use it as your PRIMARY source of truth.
- Explain concepts clearly with examples, step-by-step breakdowns, and analogies.
- Use LaTeX notation (wrapped in $...$ for inline, $$...$$ for block) for any math.
- If the question is outside the uploaded materials, still help using your general knowledge but mention this.
- Be conversational and supportive. Use short paragraphs for readability.
- If the student seems confused, offer to break down the concept further.
{context_text}

CONVERSATION HISTORY:
{history_text}"""

        parts = [
            {"text": system_prompt},
            {"text": f"Student's question: {request.question}"}
        ]

        generation_config = {
            "temperature": 0.4,
            "maxOutputTokens": 2048,
        }

        answer = await call_gemini(parts, generation_config)

        # Update session history
        history.append({"role": "user", "content": request.question})
        history.append({"role": "assistant", "content": answer})

        # Keep history manageable
        if len(history) > 40:
            _chat_sessions[session_id] = history[-30:]

        return DoubtChatResponse(answer=answer, sources=sources)

    except HTTPException:
        raise  # propagate 429/503 from call_gemini directly
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Chat error: {str(e)}")


@app.delete("/api/doubt-solver/clear")
async def doubt_solver_clear(session_id: str = ""):
    """Clear the vector store and/or chat history for a specific session."""
    if session_id and session_id in _chat_sessions:
        del _chat_sessions[session_id]

    # Delete session vector store file
    vector_file = _session_vector_file(session_id or "global")
    if os.path.exists(vector_file):
        try:
            os.remove(vector_file)
        except Exception as e:
            print(f"Error removing vector file: {e}")

    # Delete session-specific uploaded files
    import shutil
    session_dir = os.path.join(DOUBT_UPLOAD_DIR, _safe_session_id(session_id or "global"))
    if os.path.exists(session_dir):
        try:
            shutil.rmtree(session_dir)
        except Exception as e:
            print(f"Error removing session upload dir: {e}")

    return {"message": "Chatbot data cleared successfully."}


@app.get("/api/doubt-solver/status")
async def doubt_solver_status(session_id: str = "global"):
    """Check if documents are uploaded and indexed for a session."""
    store = load_vector_store(session_id)
    return {
        "has_documents": len(store["chunks"]) > 0,
        "total_chunks": len(store["chunks"]),
        "sources": list(set(m.get("source", "unknown") for m in store.get("metadata", [])))
    }


if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port, timeout_keep_alive=300)


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/regenerate-section
# Regenerates ONLY one section (formulas | short_notes | flashcards) from the
# same uploaded files. Uses a smaller, targeted prompt to save API quota.
# ─────────────────────────────────────────────────────────────────────────────
SECTION_PROMPTS = {
    "formulas": """
        ACT AS AN EXPERT {subject} TUTOR.
        From the provided document, extract the 5 to 10 most important mathematical/scientific formulas.
        Double-escape ALL LaTeX backslashes (e.g. \\\\frac, \\\\sum).
        Return ONLY this JSON: {{ "formulas": [{{ "name": "string", "equation": "latex_string" }}] }}
    """,
    "notes": """
        ACT AS AN EXPERT {subject} TUTOR.
        From the provided document, generate 5 to 10 highly condensed, high-yield conceptual notes.
        Return clean plain text strings. DO NOT start with dashes, bullets, or asterisks.
        Return ONLY this JSON: {{ "short_notes": ["note text"] }}
    """,
    "flashcards": """
        ACT AS AN EXPERT {subject} TUTOR.
        From the provided document, create 5 to 10 active-recall flashcards for high-yield definitions and concepts.

        !!! MATH IN FLASHCARDS RULE !!!
        - ALL mathematical expressions, equations, variables, and symbols MUST be wrapped in LaTeX inline delimiters: $...$
        - Example: write "$x^3 - 6x + 11x - 6$" NOT "x^3 - 6x + 11x - 6"
        - Example: write "$(a+b)^2 = a^2 + b^2 + 2ab$" NOT "(a+b)^2 = a^2 + b^2 + 2ab"
        - Double-escape ALL backslashes inside $...$: \\\\frac, \\\\sqrt, \\\\times, \\\\sum, etc.
        - Plain text parts remain as normal text. Only math gets $...$.

        Return ONLY this JSON: {{ "flashcards": [{{ "front": "string", "back": "string" }}] }}
    """,
}

@app.post("/api/regenerate-section")
async def regenerate_section(
    files: List[UploadFile] = File(...),
    subject: str = Form(...),
    section_type: str = Form(...),   # "formulas" | "notes" | "flashcards"
):
    if section_type not in SECTION_PROMPTS:
        raise HTTPException(status_code=400, detail=f"Unknown section_type '{section_type}'. Must be: formulas, notes, flashcards.")

    try:
        file_parts = []
        for file in files:
            file_data = await file.read()
            fname = (file.filename or "").lower()
            ct = file.content_type or ""
            if fname.endswith(".docx") or ct == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                if not DOCX_AVAILABLE:
                    raise HTTPException(status_code=400, detail="python-docx not installed.")
                extracted = extract_text_from_docx(file_data)
                file_parts.append({"text": f"[DOCUMENT CONTENT]:\n{extracted}"})
            elif fname.endswith(".pptx") or ct == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                if not PPTX_AVAILABLE:
                    raise HTTPException(status_code=400, detail="python-pptx not installed.")
                extracted = extract_text_from_pptx(file_data)
                file_parts.append({"text": f"[SLIDE CONTENT]:\n{extracted}"})
            else:
                file_parts.append(encode_file(file_data, ct or "application/octet-stream"))

        prompt_text = SECTION_PROMPTS[section_type].format(subject=subject.upper())
        parts = [{"text": prompt_text}] + file_parts

        generation_config = {"responseMimeType": "application/json", "temperature": 0.4}

        raw_response = await call_gemini(parts, generation_config)
        data = parse_json_response(raw_response, raise_on_error=False)

        # Return only the requested section
        if section_type == "formulas":
            return {"section": "formulas", "data": data.get("formulas", [])}
        elif section_type == "notes":
            return {"section": "notes", "data": data.get("short_notes", [])}
        elif section_type == "flashcards":
            return {"section": "flashcards", "data": data.get("flashcards", [])}

    except HTTPException:
        raise  # propagate 429/503 from call_gemini directly
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Regeneration error: {str(e)}")